"""Generate a OneNote content pack for a generic software company (Northwind Software).

Replaces the older automotive pack. Design goals:

1. Distinct, consolidated topics (no near-duplicate sections).
2. Concrete, answerable facts on every page (specific numbers, owners, values) so
   the RAG chat returns real answers instead of generic process prose.
3. Several writing styles so pages look maintained by different teams.
4. Real attachments (.md, .txt, .docx, .pptx, .pdf) on a spread of pages, each
   carrying at least one unique fact so attachment retrieval is demonstrable.

Output: ``generated_onenote_pages/`` — one folder per section, one .html per page,
plus attachment files named ``<NN>_<Title>__<attachment>.<ext>`` next to the page.

Import into OneNote: copy the rendered HTML into a new page; attach the
``__`` files to that same page; then run an OneNote bootstrap reindex.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from html import escape
from pathlib import Path
import re
import shutil

OUTPUT_DIR = Path("generated_onenote_pages")
COMPANY = "Northwind Software"


@dataclass(frozen=True, slots=True)
class AttachmentSpec:
    filename: str          # e.g. "coding-standards.md"
    kind: str              # md | txt | docx | pptx | pdf
    title: str
    intro: str
    bullets: tuple[str, ...] = ()
    sections: tuple[tuple[str, tuple[str, ...]], ...] = ()  # (heading, lines)


@dataclass(frozen=True, slots=True)
class PageSpec:
    title: str
    owner: str
    summary: str
    facts: tuple[str, ...]
    steps: tuple[str, ...]
    systems: tuple[str, ...]
    commands: tuple[str, ...] = ()         # rendered as a shell code block when present
    attachments: tuple[AttachmentSpec, ...] = ()


@dataclass(frozen=True, slots=True)
class SectionSpec:
    name: str
    page_type: str
    audience: str
    process_heading: str
    evidence: str
    pages: tuple[PageSpec, ...]


# --------------------------------------------------------------------------------------
# Reusable attachment builders
# --------------------------------------------------------------------------------------

def _md(filename: str, title: str, intro: str, bullets=(), sections=()) -> AttachmentSpec:
    return AttachmentSpec(filename, "md", title, intro, tuple(bullets), tuple(sections))


def _txt(filename: str, title: str, intro: str, bullets=()) -> AttachmentSpec:
    return AttachmentSpec(filename, "txt", title, intro, tuple(bullets))


def _docx(filename: str, title: str, intro: str, bullets=(), sections=()) -> AttachmentSpec:
    return AttachmentSpec(filename, "docx", title, intro, tuple(bullets), tuple(sections))


def _pptx(filename: str, title: str, intro: str, bullets=(), sections=()) -> AttachmentSpec:
    return AttachmentSpec(filename, "pptx", title, intro, tuple(bullets), tuple(sections))


def _pdf(filename: str, title: str, intro: str, bullets=(), sections=()) -> AttachmentSpec:
    return AttachmentSpec(filename, "pdf", title, intro, tuple(bullets), tuple(sections))


# --------------------------------------------------------------------------------------
# Content taxonomy
# --------------------------------------------------------------------------------------

SECTIONS: tuple[SectionSpec, ...] = (
    SectionSpec(
        name="Onboarding",
        page_type="onboarding guide",
        audience="new hires, managers, onboarding buddies",
        process_heading="Onboarding Steps",
        evidence="completion checklist, manager confirmation, buddy sign-off",
        pages=(
            PageSpec(
                "First Week Plan for New Engineers",
                "Engineering Operations",
                "What a new software engineer should accomplish in the first five working days.",
                (
                    "Day 1 is laptop handover and account activation; accounts are provisioned within 4 business hours of the IT ticket.",
                    "Every new engineer is paired with an onboarding buddy for the first 30 days.",
                    "By end of week one the new hire ships one small change to the internal sandbox repository.",
                    "Mandatory security awareness training must be completed within the first 5 days.",
                    "The onboarding buddy books a 30-minute daily check-in for the first week.",
                ),
                (
                    "Collect the laptop and sign the asset handover form from IT.",
                    "Activate SSO and enroll a second MFA factor before opening any tool.",
                    "Clone the sandbox repository and run the bootstrap script.",
                    "Open the first starter ticket tagged 'good-first-issue' with your buddy.",
                    "Complete security awareness training and confirm to your manager.",
                ),
                ("Okta SSO", "Service Desk", "Learning Platform", "GitHub"),
                attachments=(
                    _pptx(
                        "welcome-deck.pptx",
                        f"Welcome to {COMPANY}",
                        "Orientation deck used by the buddy on day one.",
                        sections=(
                            ("Your First Day", ("Laptop handover at 09:30", "Activate Okta SSO + MFA", "Meet your onboarding buddy")),
                            ("Who To Contact", ("People team: people@northwind.local", "Service Desk: ext. 4500", "Security: security@northwind.local")),
                            ("First Week Goal", ("Ship one change to the sandbox repo", "Finish security awareness training")),
                        ),
                    ),
                    _docx(
                        "first-week-plan.docx",
                        "First Week Plan",
                        "Detailed day-by-day plan handed to every new engineer.",
                        sections=(
                            ("Day 1", ("Laptop + accounts", "MFA enrollment", "Read the Engineering Handbook intro")),
                            ("Day 2-3", ("Local environment setup", "Pair with buddy on a starter ticket")),
                            ("Day 4-5", ("Open first pull request", "Complete security training", "Manager check-in")),
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Developer Workstation Setup",
                "Developer Platform",
                "Baseline tools every engineer installs before joining a project.",
                (
                    "The standard toolchain is Git, Docker Desktop, Node.js 20 LTS, Python 3.12, and the JetBrains or VS Code IDE.",
                    "Company laptops ship with full-disk encryption enabled; never disable it.",
                    "The internal package registry is at registry.northwind.local and requires an SSO token.",
                    "Git commits must be signed; the GPG key is registered in GitHub during setup.",
                    "WSL2 with Ubuntu 22.04 is the supported Linux environment on Windows laptops.",
                ),
                (
                    "Install Git, Docker Desktop, Node.js 20, and Python 3.12.",
                    "Authenticate to the internal registry with 'nw login'.",
                    "Generate and register a signing key for Git.",
                    "Run the workstation doctor to verify the setup.",
                ),
                ("GitHub", "Docker Desktop", "Internal Registry", "Okta SSO"),
                commands=(
                    "winget install Git.Git Docker.DockerDesktop OpenJS.NodeJS.LTS Python.Python.3.12",
                    "nw login --sso",
                    "nw doctor --check git,docker,node,python,registry",
                ),
                attachments=(
                    _md(
                        "workstation-checklist.md",
                        "Workstation Setup Checklist",
                        "Copy-paste checklist to confirm a workstation is ready.",
                        bullets=(
                            "Git installed and commit signing verified",
                            "Docker Desktop running with 8 GB memory allocated",
                            "Node.js 20 LTS and Python 3.12 on PATH",
                            "Authenticated to registry.northwind.local",
                            "`nw doctor` reports all green",
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Accounts and Access on Day One",
                "Identity Team",
                "Which accounts a new hire receives automatically and which require a request.",
                (
                    "SSO, email, Slack, and GitHub are provisioned automatically from the HR start record.",
                    "Production and customer-data access is never granted on day one and requires manager + security approval.",
                    "Access requests are submitted in the Access Portal and approved within 1 business day.",
                    "Contractors receive time-boxed accounts that expire on their contract end date.",
                ),
                (
                    "Confirm SSO, email, Slack, and GitHub work after first login.",
                    "Request team-specific tools through the Access Portal.",
                    "Ask your manager to approve any elevated access.",
                ),
                ("Okta SSO", "Access Portal", "Slack", "GitHub"),
            ),
            PageSpec(
                "Team Rituals and Cadence",
                "Engineering Operations",
                "Recurring meetings and planning cycles a new hire should expect.",
                (
                    "Teams run two-week sprints starting on Wednesdays.",
                    "Daily standup is 15 minutes at 10:00 in the team channel huddle.",
                    "Sprint review and retro happen on the second Tuesday afternoon.",
                    "Incident reviews are blameless and scheduled within 3 business days of resolution.",
                ),
                (
                    "Add the team ceremony calendar to your schedule.",
                    "Post your standup update before 10:00 if you cannot attend live.",
                    "Bring one improvement idea to the retro.",
                ),
                ("Calendar", "Slack", "Jira"),
            ),
            PageSpec(
                "First Month Milestones",
                "Engineering Operations",
                "What 'ramped up' means by the end of the first month.",
                (
                    "By day 30 the engineer has merged at least five pull requests.",
                    "By day 30 the engineer can deploy to staging without help.",
                    "A 30-day manager review confirms readiness for on-call shadowing.",
                    "On-call shadowing starts only after the engineer has completed the runbook walkthrough.",
                ),
                (
                    "Track milestones with your manager in the 1:1 document.",
                    "Schedule the 30-day review in week four.",
                    "Begin on-call shadowing after the review.",
                ),
                ("Jira", "Calendar", "Runbook Library"),
            ),
            PageSpec(
                "Onboarding Buddy Responsibilities",
                "People Operations",
                "What is expected from an assigned onboarding buddy.",
                (
                    "The buddy is a peer engineer, not the new hire's manager.",
                    "Buddies commit roughly 3 hours per week for the first two weeks.",
                    "The buddy reviews the new hire's first pull request within the same day.",
                    "Buddy assignments rotate so no engineer buddies more than once per quarter.",
                ),
                (
                    "Introduce the new hire to the team and key contacts.",
                    "Pair on the first starter ticket.",
                    "Review the first pull request promptly.",
                ),
                ("Slack", "GitHub", "HR Portal"),
            ),
        ),
    ),
    SectionSpec(
        name="Engineering Handbook",
        page_type="engineering standard",
        audience="all software engineers",
        process_heading="How We Apply It",
        evidence="linked pull request, ADR record, review note",
        pages=(
            PageSpec(
                "Git Branching and Pull Request Workflow",
                "Developer Platform",
                "The trunk-based branching model and pull request rules every repo follows.",
                (
                    "We use trunk-based development: short-lived branches merged into main, no long-running release branches.",
                    "Pull requests require one approval, and two approvals for changes touching auth or billing.",
                    "Branch names follow the pattern type/ticket-id-short-description, e.g. feat/NW-1234-add-export.",
                    "main is protected: direct pushes are blocked and CI must pass before merge.",
                    "Stale branches with no activity for 30 days are deleted automatically.",
                ),
                (
                    "Branch from main with a ticket-prefixed name.",
                    "Open a pull request early as a draft.",
                    "Get the required approvals and a green CI run.",
                    "Squash-merge into main.",
                ),
                ("GitHub", "CI Pipeline", "Jira"),
                attachments=(
                    _pdf(
                        "git-workflow-cheatsheet.pdf",
                        "Git Workflow Cheat Sheet",
                        "One-page reference for the branching and PR rules.",
                        bullets=(
                            "Trunk-based: branch from main, merge back fast",
                            "Branch name: type/NW-1234-short-description",
                            "1 approval normally, 2 for auth or billing changes",
                            "Squash-merge only; CI must be green",
                            "Branches idle 30 days are auto-deleted",
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Code Review Standards",
                "Developer Platform",
                "What reviewers check and the expected turnaround.",
                (
                    "Reviewers respond within one business day; if unavailable, reassign.",
                    "Reviews focus on correctness, tests, readability, and security, in that order.",
                    "A pull request over 400 changed lines should be split unless it is generated code.",
                    "Authors must respond to every review comment before merge.",
                    "Approving a PR means you share responsibility for the change.",
                ),
                (
                    "Read the ticket and PR description before the diff.",
                    "Check tests exist and cover the change.",
                    "Leave actionable comments, not vague ones.",
                    "Approve only when you would be comfortable owning the change.",
                ),
                ("GitHub", "CI Pipeline", "SonarQube"),
            ),
            PageSpec(
                "Coding Standards and Formatting",
                "Architecture Office",
                "Language formatters and linters enforced in CI.",
                (
                    "Python uses ruff and black; TypeScript uses ESLint and Prettier.",
                    "Formatting is enforced in CI and a failing format check blocks merge.",
                    "Maximum line length is 100 characters for Python and 100 for TypeScript.",
                    "Public functions require docstrings or JSDoc describing parameters and return values.",
                    "Type checking runs with mypy (strict) and tsc (strict) on every PR.",
                ),
                (
                    "Install the pre-commit hooks with 'pre-commit install'.",
                    "Run formatters locally before pushing.",
                    "Fix lint and type errors rather than suppressing them.",
                ),
                ("CI Pipeline", "GitHub", "pre-commit"),
                attachments=(
                    _md(
                        "coding-standards.md",
                        "Coding Standards",
                        "The normative coding standards enforced across repositories.",
                        sections=(
                            ("Python", ("Format with black, lint with ruff", "Type-check with mypy --strict", "Max line length 100")),
                            ("TypeScript", ("Format with Prettier, lint with ESLint", "tsc --strict must pass", "Prefer named exports")),
                            ("General", ("Public APIs are documented", "No commented-out code in main", "Tests live next to the code")),
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Architecture Decision Records",
                "Architecture Office",
                "How and when to write an ADR.",
                (
                    "ADRs live in docs/adr/ in each repository and are numbered sequentially.",
                    "Write an ADR for any decision that is costly to reverse, such as a database or framework choice.",
                    "An ADR has four parts: context, decision, status, and consequences.",
                    "ADRs are immutable once accepted; a new ADR supersedes an old one.",
                ),
                (
                    "Copy the ADR template into docs/adr/.",
                    "Fill in context, decision, and consequences.",
                    "Open a PR and get architecture review.",
                    "Mark the ADR accepted when merged.",
                ),
                ("GitHub", "Architecture Wiki"),
                attachments=(
                    _docx(
                        "adr-template.docx",
                        "Architecture Decision Record Template",
                        "Standard template to copy for each new decision.",
                        sections=(
                            ("Context", ("What is the problem and the forces at play?",)),
                            ("Decision", ("What did we decide, and why this option?",)),
                            ("Status", ("Proposed / Accepted / Superseded",)),
                            ("Consequences", ("What becomes easier or harder as a result?",)),
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Testing Strategy",
                "Quality Engineering",
                "The testing pyramid and coverage expectations.",
                (
                    "Target coverage is 80% lines for changed code, enforced on the diff, not the whole repo.",
                    "Unit tests run on every PR; integration tests run on merge to main.",
                    "End-to-end tests run nightly against staging.",
                    "Flaky tests are quarantined within 24 hours and fixed within a week.",
                ),
                (
                    "Write unit tests alongside the code.",
                    "Add an integration test for any new external dependency.",
                    "Keep tests deterministic and independent.",
                ),
                ("CI Pipeline", "GitHub", "Playwright"),
                attachments=(
                    _pptx(
                        "testing-pyramid.pptx",
                        "Testing Strategy",
                        "Deck used in the engineering guild talk on testing.",
                        sections=(
                            ("The Pyramid", ("Many fast unit tests", "Fewer integration tests", "A thin layer of end-to-end tests")),
                            ("Coverage", ("80% on changed lines, enforced on the diff", "Whole-repo coverage is not gated")),
                            ("Flaky Tests", ("Quarantine within 24 hours", "Fix within a week")),
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Naming Conventions",
                "Architecture Office",
                "Naming rules for repositories, services, and environment variables.",
                (
                    "Repositories use kebab-case, e.g. billing-service.",
                    "Kubernetes services use the pattern <team>-<service>, e.g. payments-billing.",
                    "Environment variables are UPPER_SNAKE_CASE and prefixed by component.",
                    "Feature flags use the pattern area.flag-name, e.g. checkout.new-cart.",
                ),
                (
                    "Check the catalog before naming something new.",
                    "Propose additions via a PR to the naming catalog.",
                ),
                ("Architecture Wiki", "GitHub"),
            ),
        ),
    ),
    SectionSpec(
        name="Project Setups",
        page_type="project setup record",
        audience="software engineers joining a project",
        process_heading="Setup Process",
        evidence="setup ticket, bootstrap output, verification log",
        pages=(
            PageSpec(
                "Billing Service Setup",
                "Payments Team",
                "Local setup for the billing-service backend (Python/FastAPI).",
                (
                    "billing-service is a FastAPI app backed by PostgreSQL and Stripe.",
                    "It runs locally on port 8020 and uses a seeded test database.",
                    "Stripe test keys are pulled from the Secrets Vault path secret/billing/stripe-test.",
                    "The smoke test creates a test invoice and asserts a paid webhook is received.",
                ),
                (
                    "Clone the repo and create a virtualenv.",
                    "Pull Stripe test secrets from the vault.",
                    "Start dependencies with docker compose.",
                    "Run migrations and the smoke test.",
                ),
                ("GitHub", "PostgreSQL", "Secrets Vault", "Stripe"),
                commands=(
                    "git clone ssh://git.northwind.local/payments/billing-service.git",
                    "cd billing-service && python -m venv .venv && . .venv/bin/activate",
                    "pip install -e '.[dev]'",
                    "nw vault read secret/billing/stripe-test > .env.local",
                    "docker compose up -d postgres",
                    "alembic upgrade head && pytest tests/smoke",
                ),
                attachments=(
                    _md(
                        "billing-service-readme.md",
                        "billing-service README",
                        "Quick reference for running billing-service locally.",
                        sections=(
                            ("Ports", ("API: 8020", "Postgres: 5432")),
                            ("Key commands", ("alembic upgrade head", "pytest tests/smoke", "uvicorn billing.main:app --port 8020")),
                            ("Gotchas", ("Stripe webhooks need the CLI: stripe listen --forward-to localhost:8020/webhooks",)),
                        ),
                    ),
                    _docx(
                        "billing-architecture-overview.docx",
                        "Billing Service Architecture Overview",
                        "How billing-service fits into the platform.",
                        sections=(
                            ("Responsibilities", ("Invoice creation", "Payment reconciliation", "Dunning emails")),
                            ("Dependencies", ("PostgreSQL for state", "Stripe for charges", "Kafka topic billing.events for downstream")),
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Web App Setup",
                "Frontend Team",
                "Local setup for the customer web app (React/TypeScript/Vite).",
                (
                    "web-app is a React 19 + TypeScript app built with Vite.",
                    "It runs on port 5173 and proxies the API gateway at api.dev.northwind.local.",
                    "Component tests use Vitest; end-to-end tests use Playwright.",
                    "Design tokens come from the shared @northwind/ui package.",
                ),
                (
                    "Clone the repo and install dependencies with npm.",
                    "Copy the example env file.",
                    "Start the dev server and run the tests.",
                ),
                ("GitHub", "Vite", "API Gateway", "Playwright"),
                commands=(
                    "git clone ssh://git.northwind.local/frontend/web-app.git",
                    "cd web-app && npm install",
                    "cp .env.example .env.local",
                    "npm run dev",
                    "npm run test && npm run e2e",
                ),
                attachments=(
                    _txt(
                        "web-app-env.txt",
                        "web-app environment variables",
                        "Example environment values for local development.",
                        bullets=(
                            "VITE_API_BASE_URL=https://api.dev.northwind.local",
                            "VITE_FEATURE_NEW_CART=true",
                            "VITE_SENTRY_DSN= (leave blank locally)",
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Search Service Setup",
                "Platform Team",
                "Local setup for the search-service (Go) backed by OpenSearch.",
                (
                    "search-service is written in Go and indexes documents into OpenSearch.",
                    "It runs on port 8040 and depends on a local OpenSearch container.",
                    "The reindex command rebuilds the index from the catalog database.",
                    "Health is exposed at /healthz and readiness at /readyz.",
                ),
                (
                    "Clone the repo and start OpenSearch with docker compose.",
                    "Build the service with make.",
                    "Run the reindex and the smoke test.",
                ),
                ("GitHub", "OpenSearch", "Docker"),
                commands=(
                    "git clone ssh://git.northwind.local/platform/search-service.git",
                    "cd search-service && docker compose up -d opensearch",
                    "make build",
                    "./bin/search-service reindex --from catalog",
                    "make smoke",
                ),
            ),
            PageSpec(
                "Data Pipeline Setup",
                "Data Team",
                "Local setup for the analytics data pipeline (Python/Airflow).",
                (
                    "The pipeline orchestrates daily ETL jobs with Apache Airflow.",
                    "Airflow runs locally on port 8080 with the example DAGs disabled.",
                    "Source data is mocked from data/fixtures during local runs.",
                    "The dbt models build into a local DuckDB file for fast iteration.",
                ),
                (
                    "Start Airflow with the provided docker compose file.",
                    "Seed local fixtures.",
                    "Trigger the daily_etl DAG and confirm success.",
                ),
                ("GitHub", "Airflow", "dbt", "DuckDB"),
                commands=(
                    "git clone ssh://git.northwind.local/data/analytics-pipeline.git",
                    "cd analytics-pipeline && docker compose up -d airflow",
                    "python scripts/seed_fixtures.py",
                    "airflow dags trigger daily_etl",
                    "dbt build --profiles-dir profiles/local",
                ),
            ),
            PageSpec(
                "Mobile App Setup",
                "Mobile Team",
                "Local setup for the Northwind mobile app (Flutter).",
                (
                    "The mobile app is built with Flutter and targets iOS and Android.",
                    "Flutter 3.24 is the supported version; mismatched versions fail CI.",
                    "The app points at the staging API by default via --dart-define.",
                    "Golden tests verify key screens render pixel-consistently.",
                ),
                (
                    "Install Flutter 3.24 and run flutter doctor.",
                    "Fetch packages and run the analyzer.",
                    "Run the app against staging and execute tests.",
                ),
                ("GitHub", "Flutter SDK", "Firebase"),
                commands=(
                    "git clone ssh://git.northwind.local/mobile/northwind-app.git",
                    "cd northwind-app && flutter pub get",
                    "flutter analyze && flutter test",
                    "flutter run --dart-define=ENV=staging",
                ),
            ),
            PageSpec(
                "Internal CLI Setup",
                "Developer Platform",
                "Installing the 'nw' internal command-line tool.",
                (
                    "The nw CLI bootstraps projects, reads vault secrets, and runs environment checks.",
                    "It is installed from the internal PyPI mirror and updated weekly.",
                    "nw doctor validates Git, Docker, registry auth, and vault access.",
                    "nw is required by most other project setup pages.",
                ),
                (
                    "Install nw from the internal mirror.",
                    "Authenticate with SSO.",
                    "Run nw doctor and resolve any red checks.",
                ),
                ("Internal Registry", "Okta SSO", "Secrets Vault"),
                commands=(
                    "pip install --index-url https://pypi.northwind.local/simple northwind-cli",
                    "nw login --sso",
                    "nw doctor",
                ),
            ),
        ),
    ),
    SectionSpec(
        name="Releases and Deployment",
        page_type="release process",
        audience="release managers, on-call engineers, team leads",
        process_heading="Deployment Process",
        evidence="release ticket, approval record, rollback owner, monitoring link",
        pages=(
            PageSpec(
                "Production Release Process",
                "Release Management",
                "How a change goes from merged to deployed in production.",
                (
                    "Production deploys happen via the CD pipeline after a green staging soak of at least 2 hours.",
                    "Deploys are blocked during the Friday 16:00 to Monday 09:00 freeze unless it is a hotfix.",
                    "Every release names a rollback owner before rollout starts.",
                    "Canary rollout sends 5% of traffic for 15 minutes before full rollout.",
                ),
                (
                    "Confirm staging soak passed and approvals are recorded.",
                    "Start the canary at 5% and watch error rate and latency.",
                    "Promote to 100% if the canary is clean.",
                    "Record the release note and monitoring link.",
                ),
                ("CD Pipeline", "Grafana", "Slack", "Jira"),
                attachments=(
                    _pdf(
                        "release-checklist.pdf",
                        "Production Release Checklist",
                        "Gate checklist completed before every production rollout.",
                        bullets=(
                            "Staging soak >= 2 hours, green",
                            "Rollback owner named",
                            "Canary 5% for 15 minutes, error rate flat",
                            "Dashboards and alerts confirmed",
                            "Release note posted in #releases",
                        ),
                    ),
                    _pptx(
                        "release-process-overview.pptx",
                        "Production Release Process",
                        "Overview deck for new release managers.",
                        sections=(
                            ("Before Rollout", ("Staging soak >= 2h", "Approvals recorded", "Rollback owner named")),
                            ("Rollout", ("Canary 5% for 15 minutes", "Watch error rate and p95", "Promote to 100% if clean")),
                            ("Freeze Window", ("No deploys Fri 16:00 - Mon 09:00", "Hotfixes are the only exception")),
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Rollback Procedure",
                "Release Management",
                "How to roll back a bad production release quickly.",
                (
                    "Rollback is a one-command redeploy of the previous known-good image tag.",
                    "The rollback target is always the image tagged last-stable.",
                    "Database migrations are backward compatible for at least one release to make rollback safe.",
                    "If a migration is not reversible, a forward-fix is used instead of rollback.",
                ),
                (
                    "Declare a rollback in the incident channel.",
                    "Run the rollback command to last-stable.",
                    "Confirm error rate returns to baseline.",
                    "Open a follow-up to forward-fix the root cause.",
                ),
                ("CD Pipeline", "Grafana", "PagerDuty"),
                commands=(
                    "nw deploy rollback --service <service> --to last-stable",
                    "nw deploy status --service <service>",
                ),
            ),
            PageSpec(
                "Feature Flag Management",
                "Platform Team",
                "How features are shipped dark and enabled gradually.",
                (
                    "Flags are managed in the Feature Flag Console and default to off in production.",
                    "Every flag has an owner and an expiry date; expired flags are reported weekly.",
                    "Risky flags are enabled by cohort: internal users, then 10%, then 100%.",
                    "A flag must be removed within 60 days of reaching 100%.",
                ),
                (
                    "Create the flag with an owner and expiry.",
                    "Enable for internal users first.",
                    "Ramp by cohort while watching metrics.",
                    "Remove the flag after full rollout.",
                ),
                ("Feature Flag Console", "Grafana", "Jira"),
            ),
            PageSpec(
                "Staging Environment Policy",
                "Platform Team",
                "What staging is for and how it stays trustworthy.",
                (
                    "Staging mirrors production configuration and is refreshed with anonymized data weekly.",
                    "Staging is shared; do not run destructive load tests without booking a window.",
                    "A change must pass staging before it is eligible for production.",
                    "Staging uses real third-party sandboxes, not mocks.",
                ),
                (
                    "Deploy your change to staging via the pipeline.",
                    "Run the relevant smoke and integration checks.",
                    "Leave staging in a clean state for the next team.",
                ),
                ("CD Pipeline", "Grafana"),
            ),
            PageSpec(
                "Release Notes and Changelog",
                "Product Engineering",
                "How customer-facing release notes are produced.",
                (
                    "Customer release notes are generated from PR titles labeled 'changelog'.",
                    "Notes are published every second Wednesday with the sprint release.",
                    "Breaking changes require a migration note and 30 days notice to customers.",
                    "Internal-only changes are excluded from the public changelog.",
                ),
                (
                    "Label customer-facing PRs with 'changelog'.",
                    "Review the generated draft before publishing.",
                    "Add migration notes for breaking changes.",
                ),
                ("GitHub", "Product Portal"),
            ),
        ),
    ),
    SectionSpec(
        name="Runbooks and Troubleshooting",
        page_type="runbook",
        audience="on-call engineers, support engineers",
        process_heading="Investigation Steps",
        evidence="incident timeline, diagnostic output, fix note, verification result",
        pages=(
            PageSpec(
                "High API Latency Runbook",
                "Platform Team",
                "Diagnosing and mitigating elevated API gateway latency.",
                (
                    "The page latency alert fires when p95 exceeds 800 ms for 5 minutes.",
                    "The most common cause is database connection pool exhaustion.",
                    "Check the Grafana 'API Overview' dashboard first, panel 'p95 by route'.",
                    "Scaling the API deployment to 6 replicas is the standard first mitigation.",
                ),
                (
                    "Acknowledge the PagerDuty alert.",
                    "Open the API Overview dashboard and find the slow route.",
                    "Check the database connection pool saturation panel.",
                    "Scale replicas or shed load, then verify p95 recovers.",
                ),
                ("Grafana", "PagerDuty", "Kubernetes"),
                commands=(
                    "kubectl scale deploy api-gateway --replicas=6 -n platform",
                    "kubectl get hpa -n platform",
                ),
                attachments=(
                    _md(
                        "api-latency-runbook.md",
                        "API Latency Runbook",
                        "Condensed steps for paging engineers at 3am.",
                        sections=(
                            ("Alert", ("p95 > 800ms for 5 min",)),
                            ("First checks", ("API Overview dashboard", "DB pool saturation", "Recent deploy in #releases")),
                            ("Mitigations", ("Scale api-gateway to 6 replicas", "Roll back last deploy if it correlates")),
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Failed Deployment Runbook",
                "Release Management",
                "What to do when a production deployment fails health checks.",
                (
                    "The pipeline automatically halts rollout if readiness checks fail for 3 minutes.",
                    "A failed deploy never takes down the running version; old pods keep serving.",
                    "Most failures are missing config or a failed migration.",
                    "If unsure, roll back to last-stable rather than debugging in production.",
                ),
                (
                    "Read the pipeline logs for the failing step.",
                    "Check for missing config or migration errors.",
                    "Roll back to last-stable if the fix is not obvious.",
                    "Open an incident if customers were affected.",
                ),
                ("CD Pipeline", "Grafana", "PagerDuty"),
            ),
            PageSpec(
                "Database Connection Errors",
                "Data Team",
                "Resolving 'too many connections' and timeout errors from PostgreSQL.",
                (
                    "The primary database max_connections is 400; pgbouncer pools at 1000 client connections.",
                    "A connection leak usually shows as steadily rising idle-in-transaction sessions.",
                    "Restarting the leaking service clears leaked connections immediately.",
                    "Long-running queries over 30 seconds are killed by a watchdog.",
                ),
                (
                    "Query pg_stat_activity for idle-in-transaction sessions.",
                    "Identify the owning service.",
                    "Restart that service to clear leaks.",
                    "File a bug to fix the leak at the source.",
                ),
                ("PostgreSQL", "pgbouncer", "Grafana"),
                commands=(
                    "SELECT state, count(*) FROM pg_stat_activity GROUP BY state;",
                    "SELECT pid, now()-query_start AS age, query FROM pg_stat_activity WHERE state='active' ORDER BY age DESC LIMIT 10;",
                ),
            ),
            PageSpec(
                "Docker Compose Port Conflict",
                "Developer Platform",
                "Fixing local 'port is already allocated' errors.",
                (
                    "The error means another process or a stale container holds the port.",
                    "Common conflicts are 5432 (Postgres), 6379 (Redis), and 8080 (app).",
                    "Removing orphaned containers usually frees the port.",
                    "Changing the host port mapping in an override file avoids the conflict entirely.",
                ),
                (
                    "Find what holds the port.",
                    "Stop the stale container or process.",
                    "Re-run docker compose up.",
                ),
                ("Docker Desktop", "Developer CLI"),
                commands=(
                    "docker compose down --remove-orphans",
                    "netstat -ano | findstr :5432",
                    "docker ps -a --filter publish=5432",
                ),
                attachments=(
                    _txt(
                        "common-local-ports.txt",
                        "Common Local Ports",
                        "Default host ports used by local services, to spot conflicts fast.",
                        bullets=(
                            "5432 - PostgreSQL",
                            "6379 - Redis",
                            "6333 - Qdrant",
                            "8020 - billing-service",
                            "8040 - search-service",
                            "5173 - web-app dev server",
                            "8080 - Airflow / generic app",
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Vector Search Returns No Results",
                "Platform Team",
                "Why semantic search may return nothing and how to fix it.",
                (
                    "An empty result set usually means the embedding model or vector dimension changed without a reindex.",
                    "The collection vector size must match the embedding model dimension exactly.",
                    "Recreating the collection requires a full reindex of all documents.",
                    "Check the collection info endpoint to confirm dimension and point count.",
                ),
                (
                    "Check the collection dimension and point count.",
                    "Compare against the configured embedding model dimension.",
                    "Recreate the collection and reindex if they differ.",
                ),
                ("Qdrant", "Sync Worker"),
                commands=(
                    "curl http://localhost:6333/collections/documents",
                ),
            ),
        ),
    ),
    SectionSpec(
        name="Internal Tools and Access",
        page_type="tool access guide",
        audience="engineers and team leads",
        process_heading="Request Flow",
        evidence="access ticket, owner approval, role mapping",
        pages=(
            PageSpec(
                "GitHub Access and Teams",
                "Developer Platform",
                "How repository access maps to GitHub teams.",
                (
                    "Access is granted through GitHub teams, never to individuals directly.",
                    "Each repository has a CODEOWNERS file that defines required reviewers.",
                    "Write access to infra repositories requires security approval.",
                    "Outside collaborators are reviewed every quarter.",
                ),
                (
                    "Request to join the relevant GitHub team in the Access Portal.",
                    "Wait for the team owner to approve.",
                    "Verify you can clone and open a draft PR.",
                ),
                ("GitHub", "Access Portal", "Okta SSO"),
            ),
            PageSpec(
                "Secrets Vault Usage",
                "Security Engineering",
                "Reading and rotating secrets safely.",
                (
                    "Secrets live in HashiCorp Vault under per-team paths like secret/<team>/<app>.",
                    "Secrets are never committed to Git; CI reads them at deploy time.",
                    "Production secrets rotate every 90 days automatically.",
                    "Reading a production secret is audited and may trigger a review.",
                ),
                (
                    "Authenticate to Vault with your SSO token.",
                    "Read only the path your service needs.",
                    "Never paste secrets into chat or tickets.",
                ),
                ("Secrets Vault", "Okta SSO", "CI Pipeline"),
                commands=(
                    "nw vault login --sso",
                    "nw vault read secret/<team>/<app>",
                ),
            ),
            PageSpec(
                "CI Pipeline Access and Reruns",
                "Build Infrastructure",
                "How to view, rerun, and debug CI jobs.",
                (
                    "CI runs on GitHub Actions with self-hosted runners labeled nw-linux and nw-macos.",
                    "Anyone on the repo can rerun a failed job; only owners can edit workflows.",
                    "A job stuck in queue for over 10 minutes usually means no runner is available.",
                    "Build artifacts are retained for 14 days.",
                ),
                (
                    "Open the Actions tab on the repository.",
                    "Rerun the failed job or download logs.",
                    "Escalate to Build Infrastructure if no runner picks it up.",
                ),
                ("GitHub", "CI Pipeline", "Service Desk"),
            ),
            PageSpec(
                "Observability Stack Access",
                "Platform Team",
                "Getting into Grafana, logs, and traces.",
                (
                    "Dashboards are in Grafana at grafana.northwind.local with SSO.",
                    "Logs are in Loki and queried through Grafana Explore.",
                    "Traces are in Tempo, linked from each service dashboard.",
                    "Alert routing is configured per team in the Alerting section.",
                ),
                (
                    "Sign in to Grafana with SSO.",
                    "Find your team folder of dashboards.",
                    "Use Explore for ad-hoc log and trace queries.",
                ),
                ("Grafana", "Loki", "Tempo", "Okta SSO"),
            ),
            PageSpec(
                "Jira Projects and Boards",
                "Engineering Operations",
                "Finding the right board and requesting access.",
                (
                    "Each team has a Jira project keyed by a short code, e.g. NW for the platform team.",
                    "Tickets follow the workflow Backlog, In Progress, In Review, Done.",
                    "Access is requested through the Access Portal and approved by the project lead.",
                    "Cross-team work uses linked issues rather than moving tickets between projects.",
                ),
                (
                    "Request access to the project in the Access Portal.",
                    "Use the team board filter for your sprint.",
                    "Link cross-team dependencies instead of copying tickets.",
                ),
                ("Jira", "Access Portal"),
            ),
        ),
    ),
    SectionSpec(
        name="IT Support",
        page_type="IT support procedure",
        audience="employees and service desk agents",
        process_heading="Support Process",
        evidence="ticket ID, requester identity, resolution note",
        pages=(
            PageSpec(
                "VPN Access and Recovery",
                "Network Services",
                "Connecting to the corporate VPN and recovering lost access.",
                (
                    "The VPN client is Cisco AnyConnect pointing at vpn.northwind.local.",
                    "VPN requires SSO plus a push approval on your enrolled MFA device.",
                    "If the certificate is invalid, reinstall the profile from the Self-Service Portal.",
                    "VPN is only required for admin networks; most SaaS tools work without it.",
                ),
                (
                    "Open AnyConnect and connect to vpn.northwind.local.",
                    "Approve the MFA push.",
                    "If it fails, reinstall the profile from Self-Service.",
                ),
                ("VPN Gateway", "Okta SSO", "Self-Service Portal"),
            ),
            PageSpec(
                "MFA Device Replacement",
                "Identity Team",
                "Re-enrolling MFA after getting a new phone.",
                (
                    "You must verify identity with a manager video confirmation before MFA reset.",
                    "Reset is performed by the Service Desk, not self-service, to prevent account takeover.",
                    "After reset you have 24 hours to enroll the new device before access is suspended.",
                    "Hardware security keys are recommended for engineers with production access.",
                ),
                (
                    "Open a Service Desk ticket for MFA reset.",
                    "Complete the manager identity verification.",
                    "Enroll your new device within 24 hours.",
                ),
                ("Okta SSO", "Service Desk", "HR Portal"),
            ),
            PageSpec(
                "Laptop Replacement and Repair",
                "Endpoint Engineering",
                "Getting a broken or lost laptop replaced.",
                (
                    "Standard replacement laptops ship within 2 business days from the local depot.",
                    "A lost laptop must be reported within 1 hour so it can be remotely wiped.",
                    "All laptops are encrypted, so a lost device does not expose data if reported.",
                    "Your data is recoverable if you used the synced Documents folder.",
                ),
                (
                    "Report the issue to the Service Desk.",
                    "For a lost device, trigger remote wipe immediately.",
                    "Collect the replacement and restore from sync.",
                ),
                ("Service Desk", "MDM Portal", "Asset Register"),
            ),
            PageSpec(
                "Software Installation Requests",
                "Endpoint Engineering",
                "Installing software outside the standard image.",
                (
                    "Pre-approved software installs instantly from the Self-Service Portal.",
                    "Non-standard software needs a license owner and a security review.",
                    "Engineers can install developer tools from the approved Homebrew/winget list without a ticket.",
                    "Unapproved software is blocked by endpoint policy.",
                ),
                (
                    "Check the Self-Service Portal first.",
                    "For other software, open a request with a business reason.",
                    "Provide license details for paid software.",
                ),
                ("Self-Service Portal", "Service Desk", "Security Review"),
            ),
            PageSpec(
                "Email and Distribution Lists",
                "Collaboration Services",
                "Requesting shared mailboxes and distribution lists.",
                (
                    "Distribution lists follow the pattern team-<name>@northwind.local.",
                    "Every list needs a named owner who approves membership.",
                    "Shared mailboxes are reviewed annually for active owners.",
                    "External senders to internal lists are quarantined by default.",
                ),
                (
                    "Request the list in the Access Portal with an owner.",
                    "Define the membership rule.",
                    "Confirm delivery with a test message.",
                ),
                ("Mail Admin Console", "Access Portal"),
            ),
        ),
    ),
    SectionSpec(
        name="People and HR",
        page_type="HR policy",
        audience="employees, managers, HR partners",
        process_heading="Employee Process",
        evidence="policy source, request record, approval",
        pages=(
            PageSpec(
                "Paid Time Off Policy",
                "People Operations",
                "How annual leave, carryover, and approval work.",
                (
                    "Full-time employees receive 25 days of annual leave per calendar year.",
                    "Up to 5 unused days carry over and expire on 31 March of the next year.",
                    "Leave is requested in the HR Portal and approved by the direct manager.",
                    "Requests of 5 or more consecutive days should be submitted 3 weeks in advance.",
                ),
                (
                    "Submit the leave request in the HR Portal.",
                    "Get manager approval.",
                    "Arrange coverage for any on-call duty.",
                ),
                ("HR Portal", "Calendar"),
                attachments=(
                    _pdf(
                        "pto-policy-summary.pdf",
                        "PTO Policy Summary",
                        "One-page summary of the paid time off rules.",
                        bullets=(
                            "25 days annual leave per year (full-time)",
                            "Up to 5 days carry over, expire 31 March",
                            "Request via HR Portal, manager approves",
                            "5+ consecutive days: 3 weeks notice",
                            "Sick leave is separate and not capped by PTO",
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Remote and Hybrid Work",
                "People Operations",
                "Expectations for where and when employees work.",
                (
                    "The default is hybrid: at least 2 days per week in the office for office-assigned roles.",
                    "Fully remote arrangements require VP approval and a signed remote work agreement.",
                    "Core collaboration hours are 10:00 to 16:00 in the employee's home time zone.",
                    "Working from abroad over 30 days requires a tax and security review.",
                ),
                (
                    "Agree your pattern with your manager.",
                    "For fully remote, request VP approval.",
                    "For long stays abroad, start the tax review early.",
                ),
                ("HR Portal", "Identity Provider"),
            ),
            PageSpec(
                "On-Call Compensation",
                "People Operations",
                "How engineers are compensated for on-call duty.",
                (
                    "On-call is one week at a time and rotates so no one does more than one week per month.",
                    "On-call carries a flat weekly stipend plus time off in lieu for night pages.",
                    "Being paged after midnight grants a late start the next day.",
                    "On-call eligibility starts only after on-call shadowing is complete.",
                ),
                (
                    "Confirm the rotation in the schedule.",
                    "Record any night pages for time in lieu.",
                    "Hand off cleanly at the end of the week.",
                ),
                ("PagerDuty", "HR Portal", "Payroll System"),
            ),
            PageSpec(
                "Performance Review Cycle",
                "People Operations",
                "Timing and inputs for performance reviews.",
                (
                    "Reviews run twice a year, in June and December.",
                    "Each review includes a self-assessment and at least two peer reviews.",
                    "Calibration happens at the department level before ratings are shared.",
                    "Promotion cases are submitted during the December cycle.",
                ),
                (
                    "Write your self-assessment by the deadline.",
                    "Nominate peer reviewers.",
                    "Discuss outcomes in your 1:1.",
                ),
                ("HR Portal", "Performance Tool"),
            ),
            PageSpec(
                "Health Insurance and Benefits Enrollment",
                "Benefits Team",
                "Enrolling in health insurance and core benefits.",
                (
                    "New hires have 30 days from their start date to enroll in health insurance.",
                    "The company covers 100% of the employee premium and 60% for dependents.",
                    "Open enrollment for changes runs every November.",
                    "A qualifying life event allows changes outside open enrollment within 30 days.",
                ),
                (
                    "Open the Benefits Portal within 30 days of starting.",
                    "Select your plan and add dependents.",
                    "Upload any required dependent proof.",
                ),
                ("Benefits Portal", "HR Portal", "Insurance Provider"),
                attachments=(
                    _docx(
                        "benefits-overview.docx",
                        "Benefits Overview",
                        "Summary of the benefits package for new hires.",
                        sections=(
                            ("Health", ("100% employee premium covered", "60% dependent premium covered", "Dental and vision included")),
                            ("Wellness", ("Annual wellness budget of 500 EUR", "Employee assistance program")),
                            ("Retirement", ("Employer pension match up to 5%",)),
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Learning and Development Budget",
                "Learning and Development",
                "Using the annual training budget.",
                (
                    "Each engineer has an annual learning budget of 1,500 EUR.",
                    "Conferences, courses, books, and certifications are eligible.",
                    "Manager approval is required before booking; reimbursement needs a receipt.",
                    "Unused budget does not carry over to the next year.",
                ),
                (
                    "Pick an eligible course or conference.",
                    "Get manager approval in the HR Portal.",
                    "Submit the receipt for reimbursement.",
                ),
                ("HR Portal", "Learning Platform", "Expense Tool"),
            ),
        ),
    ),
    SectionSpec(
        name="Finance and Procurement",
        page_type="finance process",
        audience="managers, project owners, finance partners",
        process_heading="Approval Process",
        evidence="approval record, cost center, receipt or quote",
        pages=(
            PageSpec(
                "Expense Reimbursement",
                "Finance Operations",
                "How to claim work-related expenses.",
                (
                    "Expenses are submitted in the Expense Tool within 60 days of the purchase.",
                    "Receipts are mandatory for any expense over 25 EUR.",
                    "Approved expenses are paid with the next monthly payroll run.",
                    "Alcohol and personal items are not reimbursable.",
                ),
                (
                    "Photograph the receipt and create an expense.",
                    "Assign the correct cost center.",
                    "Submit for manager approval.",
                ),
                ("Expense Tool", "Finance Approval", "Payroll System"),
                attachments=(
                    _pdf(
                        "expense-policy.pdf",
                        "Expense Policy",
                        "What is and is not reimbursable.",
                        bullets=(
                            "Submit within 60 days",
                            "Receipts required over 25 EUR",
                            "Paid with next monthly payroll",
                            "No alcohol or personal items",
                            "Travel booked through the corporate portal",
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Software and SaaS Purchasing",
                "Procurement",
                "Buying new software tools and subscriptions.",
                (
                    "Any new SaaS tool requires a security and data-privacy review before purchase.",
                    "Purchases over 5,000 EUR per year need finance director approval.",
                    "Prefer existing approved vendors before evaluating new ones.",
                    "Annual contracts are preferred over monthly for tools used by whole teams.",
                ),
                (
                    "Check the approved vendor list first.",
                    "Submit a purchase request with a business case.",
                    "Complete the security review for new vendors.",
                ),
                ("Procurement Portal", "Security Review", "Finance Approval"),
            ),
            PageSpec(
                "Cloud Cost Management",
                "Cloud FinOps",
                "How teams stay accountable for cloud spend.",
                (
                    "Every cloud resource must carry team, service, and environment tags.",
                    "Untagged resources are flagged daily and may be stopped after 7 days.",
                    "Each team gets a monthly cost report and a budget alert at 80% of plan.",
                    "Non-production environments scale down outside working hours.",
                ),
                (
                    "Tag all resources at creation.",
                    "Review the monthly cost report with your team.",
                    "Investigate any spend anomaly alert.",
                ),
                ("Cloud Console", "FinOps Dashboard"),
            ),
            PageSpec(
                "Vendor Invoice Processing",
                "Accounts Payable",
                "How supplier invoices are received and paid.",
                (
                    "Invoices are emailed to invoices@northwind.local or uploaded to the Invoice Portal.",
                    "An invoice must match an approved purchase order to be paid.",
                    "Standard payment terms are net 30 from invoice date.",
                    "Disputed invoices are put on hold and the supplier is notified within 5 days.",
                ),
                (
                    "Match the invoice to its purchase order.",
                    "Confirm goods or services were received.",
                    "Approve for payment or raise a dispute.",
                ),
                ("Invoice Portal", "ERP", "Procurement Portal"),
            ),
            PageSpec(
                "Team Budget and Forecast",
                "Program Finance",
                "How engineering managers track their budget.",
                (
                    "Budgets are set annually and reviewed each quarter.",
                    "Headcount is the largest line and is forecast in the Planning Tool.",
                    "A variance over 10% against plan requires a written explanation.",
                    "Capital purchases are depreciated and tracked separately from operating cost.",
                ),
                (
                    "Review the monthly actuals against plan.",
                    "Update the forecast in the Planning Tool.",
                    "Explain variances over 10%.",
                ),
                ("ERP", "Planning Tool", "Finance Dashboard"),
            ),
        ),
    ),
    SectionSpec(
        name="Security and Compliance",
        page_type="security policy",
        audience="all employees, engineers, auditors",
        process_heading="Control Process",
        evidence="control ID, review record, evidence link",
        pages=(
            PageSpec(
                "Data Classification Policy",
                "Security Engineering",
                "The four data classes and how to handle each.",
                (
                    "Data is classified as Public, Internal, Confidential, or Restricted.",
                    "Customer personal data is always at least Confidential.",
                    "Restricted data, such as payment card data, may never be stored in logs or screenshots.",
                    "Confidential data must be encrypted at rest and in transit.",
                ),
                (
                    "Determine the highest class of data you handle.",
                    "Apply the matching storage and sharing rules.",
                    "Ask Security if you are unsure.",
                ),
                ("Security Portal", "Data Catalog"),
                attachments=(
                    _pdf(
                        "data-classification-policy.pdf",
                        "Data Classification Policy",
                        "Reference for the four data classes.",
                        sections=(
                            ("Public", ("Marketing site, open-source code",)),
                            ("Internal", ("Most internal docs and tickets",)),
                            ("Confidential", ("Customer personal data, contracts; encrypt at rest",)),
                            ("Restricted", ("Payment card data, secrets; never in logs",)),
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Security Incident Response",
                "Security Operations",
                "What to do when you suspect a security incident.",
                (
                    "Report suspected incidents within 30 minutes to security@northwind.local or #security-incident.",
                    "Do not investigate alone; the on-call security engineer coordinates response.",
                    "Preserve evidence: do not delete logs or wipe affected machines.",
                    "Customer-impacting breaches trigger the legal and communications playbook.",
                ),
                (
                    "Report immediately through the security channel.",
                    "Preserve all evidence.",
                    "Follow the incident commander's instructions.",
                ),
                ("Security Portal", "PagerDuty", "Legal Intake"),
                attachments=(
                    _docx(
                        "incident-response-plan.docx",
                        "Security Incident Response Plan",
                        "Roles and steps for handling a security incident.",
                        sections=(
                            ("Roles", ("Incident Commander", "Communications Lead", "Scribe")),
                            ("Phases", ("Detect and report", "Contain", "Eradicate", "Recover", "Post-incident review")),
                            ("Timelines", ("Report within 30 min", "Status update every 60 min")),
                        ),
                    ),
                ),
            ),
            PageSpec(
                "Access Reviews",
                "Security Engineering",
                "Quarterly review of who can access what.",
                (
                    "Access to production and customer data is reviewed every quarter.",
                    "Managers must confirm or revoke each team member's access within 10 business days.",
                    "Access not confirmed in time is automatically revoked.",
                    "Service accounts are reviewed at the same cadence as people.",
                ),
                (
                    "Open the access review task when notified.",
                    "Confirm or revoke each entry.",
                    "Complete within 10 business days.",
                ),
                ("Access Portal", "Identity Provider", "Security Portal"),
            ),
            PageSpec(
                "Secure Coding Guidelines",
                "Security Engineering",
                "Baseline practices to avoid common vulnerabilities.",
                (
                    "Never build SQL with string concatenation; use parameterized queries.",
                    "All input crossing a trust boundary is validated and encoded on output.",
                    "Dependencies are scanned on every PR and high-severity findings block merge.",
                    "Secrets are read from the vault, never hardcoded or committed.",
                ),
                (
                    "Use parameterized queries and safe APIs.",
                    "Validate and encode untrusted input.",
                    "Resolve high-severity dependency alerts before merge.",
                ),
                ("GitHub", "Dependency Scanner", "Secrets Vault"),
            ),
            PageSpec(
                "Vendor Risk and SOC 2 Evidence",
                "Compliance",
                "How vendor reviews and audit evidence are organized.",
                (
                    "New vendors handling customer data must provide a current SOC 2 Type II report.",
                    "Evidence for our own SOC 2 audit is collected continuously, not at year end.",
                    "Each control has a named owner and a quarterly evidence refresh.",
                    "Auditor access is read-only and time-boxed to the audit window.",
                ),
                (
                    "Request the vendor's SOC 2 report.",
                    "File it in the vendor risk register.",
                    "Refresh control evidence each quarter.",
                ),
                ("Compliance Portal", "Document Library"),
            ),
        ),
    ),
)


SECTION_GUIDANCE = {
    "Onboarding": "For new-hire orientation, first-week and first-month tasks, accounts, and buddy expectations.",
    "Engineering Handbook": "For engineering standards: branching, code review, formatting, ADRs, testing, and naming.",
    "Project Setups": (
        "For project setup inventory and local setup instructions. Each page title in this section is one project setup. "
        "Do not treat tools, runbooks, or deployment pages as project setups."
    ),
    "Releases and Deployment": "For release process, rollout, rollback, feature flags, staging, and release notes.",
    "Runbooks and Troubleshooting": "For symptoms, diagnostics, mitigations, and verification of recurring failures.",
    "Internal Tools and Access": "For access to GitHub, Vault, CI, observability, and Jira.",
    "IT Support": "For helpdesk procedures: VPN, MFA, laptops, software installs, and email.",
    "People and HR": "For leave, remote work, on-call pay, performance reviews, health insurance, and learning budget.",
    "Finance and Procurement": "For expenses, software purchasing, cloud cost, invoices, and team budgets.",
    "Security and Compliance": "For data classification, incident response, access reviews, secure coding, and audits.",
}


HTML_STYLE = """
body { font-family: Aptos, Arial, sans-serif; color: #1f2933; line-height: 1.45; max-width: 900px; margin: 32px auto; padding: 0 24px; }
h1 { font-size: 28px; margin: 0 0 8px; }
h2 { font-size: 18px; margin-top: 28px; border-bottom: 1px solid #d9e2ec; padding-bottom: 4px; }
p.summary { font-size: 15px; margin-top: 0; color: #334e68; }
table { border-collapse: collapse; width: 100%; margin: 12px 0; }
th, td { border: 1px solid #bcccdc; padding: 7px 9px; text-align: left; vertical-align: top; }
th { background: #f0f4f8; width: 210px; }
dl { display: grid; grid-template-columns: 180px 1fr; gap: 6px 14px; margin: 12px 0; }
dt { font-weight: 700; color: #334e68; }
dd { margin: 0; }
ul, ol { padding-left: 24px; }
li { margin: 5px 0; }
pre { background: #102a43; color: #f0f4f8; padding: 12px; border-radius: 6px; overflow-x: auto; }
code { font-family: Consolas, "Courier New", monospace; }
.note { background: #fffbea; border-left: 4px solid #f0b429; padding: 10px 12px; }
.attach { background: #ebf8ff; border-left: 4px solid #2b6cb0; padding: 10px 12px; }
""".strip()


# --------------------------------------------------------------------------------------
# Rendering
# --------------------------------------------------------------------------------------

def main() -> None:
    if OUTPUT_DIR.exists():
        shutil.rmtree(OUTPUT_DIR)
    OUTPUT_DIR.mkdir(parents=True)
    write_readme()
    write_manifest()
    for index, section in enumerate(SECTIONS, start=1):
        folder = OUTPUT_DIR / f"{index:02d}_{safe_name(section.name)}"
        folder.mkdir()
        for page_index, page in enumerate(section.pages, start=1):
            prefix = f"{page_index:02d}_{safe_name(page.title)}"
            (folder / f"{prefix}.html").write_text(
                render_page(section, page, page_index), encoding="utf-8"
            )
            page_context = (
                f"Attached to the \"{page.title}\" {section.page_type} "
                f"in the {section.name} section, owned by {page.owner}."
            )
            for attachment in page.attachments:
                write_attachment(
                    folder / f"{prefix}__{attachment.filename}", attachment, context=page_context
                )


def render_page(section: SectionSpec, page: PageSpec, page_index: int) -> str:
    variant = (page_index + len(section.name)) % 4
    body = render_body(section, page, variant)
    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <title>{escape(page.title)}</title>
  <style>{HTML_STYLE}</style>
</head>
<body>
  <article>
    <h1>{escape(page.title)}</h1>
    <p class="summary">{escape(page.summary)}</p>
{body}
  </article>
</body>
</html>
"""


def render_body(section: SectionSpec, page: PageSpec, variant: int) -> str:
    # A clear, self-contained intro: who owns it, who it is for, what it covers,
    # and the single most important fact stated up front. The earlier template
    # produced ungrammatical leads ("This HR policy how annual leave ... work")
    # that added noise to the indexed text; this reads as a real page opener.
    lead_fact = f" {page.facts[0]}" if page.facts else ""
    overview = (
        f"{page.title} is the {section.page_type} maintained by {page.owner} "
        f"for {section.audience}. It covers {page.summary[0].lower() + page.summary[1:]}"
        f"{lead_fact}"
    )
    ownership = ownership_table(section, page) if variant % 2 == 0 else ownership_dl(section, page)
    key_facts = list_section("Key Facts", page.facts)
    steps = ordered_section(section.process_heading, page.steps)
    systems_line = paragraph_section(
        "Systems Involved",
        f"This work touches {comma_list(page.systems)}.",
    )
    commands = command_block(page.commands) if page.commands else ""
    attachments = attachment_section(page.attachments) if page.attachments else ""
    evidence = paragraph_section(
        "What To Keep",
        f"Keep the {section.evidence} with the related ticket or record.",
    )

    if variant == 0:
        parts = [paragraph_section("Overview", overview), ownership, key_facts, steps, commands, attachments, evidence]
    elif variant == 1:
        parts = [note_section("Summary", overview), key_facts, ownership, steps, commands, attachments, evidence]
    elif variant == 2:
        parts = [paragraph_section("Purpose", overview), key_facts, steps, ownership, systems_line, commands, attachments, evidence]
    else:
        parts = [paragraph_section("Context", overview), ownership, steps, key_facts, commands, attachments, systems_line]
    return "\n\n".join(part for part in parts if part.strip())


def paragraph_section(title: str, text: str) -> str:
    return f"""    <h2>{escape(title)}</h2>
    <p>{escape(text)}</p>"""


def note_section(title: str, text: str) -> str:
    return f"""    <h2>{escape(title)}</h2>
    <p class="note">{escape(text)}</p>"""


def list_section(title: str, items) -> str:
    return f"""    <h2>{escape(title)}</h2>
    <ul>
      {render_list(items)}
    </ul>"""


def ordered_section(title: str, items) -> str:
    return f"""    <h2>{escape(title)}</h2>
    <ol>
      {render_list(items)}
    </ol>"""


def ownership_table(section: SectionSpec, page: PageSpec) -> str:
    return f"""    <h2>Ownership</h2>
    <table>
      <tr><th>Owner team</th><td>{escape(page.owner)}</td></tr>
      <tr><th>Audience</th><td>{escape(section.audience)}</td></tr>
      <tr><th>Systems</th><td>{escape(comma_list(page.systems))}</td></tr>
    </table>"""


def ownership_dl(section: SectionSpec, page: PageSpec) -> str:
    return f"""    <h2>Owner and Systems</h2>
    <dl>
      <dt>Owner</dt><dd>{escape(page.owner)}</dd>
      <dt>For</dt><dd>{escape(section.audience)}</dd>
      <dt>Systems</dt><dd>{escape(comma_list(page.systems))}</dd>
    </dl>"""


def command_block(commands) -> str:
    body = "\n".join(commands)
    return f"""    <h2>Commands</h2>
    <p class="note">Run from a clean shell and keep the output with the setup or incident record.</p>
    <pre><code class="language-bash">{escape(body)}</code></pre>"""


def attachment_section(attachments) -> str:
    items = []
    for attachment in attachments:
        items.append(
            f"<li><strong>{escape(attachment.filename)}</strong> &mdash; {escape(attachment.intro)} "
            f"<em>(attach this file to the OneNote page)</em></li>"
        )
    listed = "\n      ".join(items)
    return f"""    <h2>Attachments</h2>
    <p class="attach">The files below are generated next to this page. Attach them to the OneNote page so their content is indexed and searchable.</p>
    <ul>
      {listed}
    </ul>"""


# --------------------------------------------------------------------------------------
# Attachment file writers
# --------------------------------------------------------------------------------------

def write_attachment(path: Path, attachment: AttachmentSpec, *, context: str = "") -> None:
    try:
        if attachment.kind == "md":
            path.write_text(render_markdown(attachment, context), encoding="utf-8")
        elif attachment.kind == "txt":
            path.write_text(render_plaintext(attachment, context), encoding="utf-8")
        elif attachment.kind == "docx":
            write_docx(path, attachment, context)
        elif attachment.kind == "pptx":
            write_pptx(path, attachment, context)
        elif attachment.kind == "pdf":
            write_pdf(path, attachment, context)
        else:
            raise ValueError(f"Unknown attachment kind: {attachment.kind}")
    except Exception as error:  # graceful fallback so the pack always generates
        fallback = path.with_suffix(".md")
        fallback.write_text(
            render_markdown(attachment, context) + f"\n\n> Note: original {attachment.kind} generation failed ({error}).\n",
            encoding="utf-8",
        )


def render_markdown(attachment: AttachmentSpec, context: str = "") -> str:
    lines = [f"# {attachment.title}", ""]
    if context:
        lines += [f"_{context}_", ""]
    lines += [attachment.intro, ""]
    for bullet in attachment.bullets:
        lines.append(f"- {bullet}")
    if attachment.bullets:
        lines.append("")
    for heading, body in attachment.sections:
        lines.append(f"## {heading}")
        lines.append("")
        for item in body:
            lines.append(f"- {item}")
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def render_plaintext(attachment: AttachmentSpec, context: str = "") -> str:
    lines = [attachment.title, "=" * len(attachment.title), ""]
    if context:
        lines += [context, ""]
    lines += [attachment.intro, ""]
    for bullet in attachment.bullets:
        lines.append(f"* {bullet}")
    for heading, body in attachment.sections:
        lines.append("")
        lines.append(heading)
        lines.append("-" * len(heading))
        for item in body:
            lines.append(f"* {item}")
    return "\n".join(lines).rstrip() + "\n"


def write_docx(path: Path, attachment: AttachmentSpec, context: str = "") -> None:
    from docx import Document

    document = Document()
    document.add_heading(attachment.title, level=0)
    if context:
        document.add_paragraph(context, style="Intense Quote")
    document.add_paragraph(attachment.intro)
    for bullet in attachment.bullets:
        document.add_paragraph(bullet, style="List Bullet")
    for heading, body in attachment.sections:
        document.add_heading(heading, level=1)
        for item in body:
            document.add_paragraph(item, style="List Bullet")
    document.save(str(path))


def write_pptx(path: Path, attachment: AttachmentSpec, context: str = "") -> None:
    from pptx import Presentation
    from pptx.util import Pt

    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    bullet_layout = presentation.slide_layouts[1]

    title_slide = presentation.slides.add_slide(title_layout)
    title_slide.shapes.title.text = attachment.title
    title_slide.placeholders[1].text = f"{context}\n{attachment.intro}" if context else attachment.intro

    sections = attachment.sections or (("Highlights", attachment.bullets),)
    for heading, body in sections:
        slide = presentation.slides.add_slide(bullet_layout)
        slide.shapes.title.text = heading
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for index, item in enumerate(body):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = item
            paragraph.font.size = Pt(18)
    presentation.save(str(path))


def write_pdf(path: Path, attachment: AttachmentSpec, context: str = "") -> None:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    def line(text: str, *, height: float = 6.0) -> None:
        # new_x=LMARGIN returns the cursor to the left margin; without it a
        # width-0 multi_cell leaves x at the right edge and the next call raises
        # "Not enough horizontal space to render a single character".
        pdf.multi_cell(0, height, _latin(text), new_x="LMARGIN", new_y="NEXT")

    pdf.set_font("Helvetica", "B", 16)
    line(attachment.title, height=9)
    pdf.ln(2)
    if context:
        pdf.set_font("Helvetica", "I", 10)
        line(context)
        pdf.ln(1)
    pdf.set_font("Helvetica", size=11)
    line(attachment.intro)
    pdf.ln(2)
    for bullet in attachment.bullets:
        line(f"- {bullet}")
    for heading, body in attachment.sections:
        pdf.ln(2)
        pdf.set_font("Helvetica", "B", 12)
        line(heading, height=7)
        pdf.set_font("Helvetica", size=11)
        for item in body:
            line(f"- {item}")
    pdf.output(str(path))


def _latin(text: str) -> str:
    # fpdf2 core fonts are latin-1; replace the few non-latin glyphs we use.
    return text.replace("—", "-").replace("–", "-").encode("latin-1", "replace").decode("latin-1")


# --------------------------------------------------------------------------------------
# Manifest + README
# --------------------------------------------------------------------------------------

def write_manifest() -> None:
    rows = ["section,page_title,owner,page_type,attachments"]
    for section in SECTIONS:
        for page in section.pages:
            attachments = "; ".join(a.filename for a in page.attachments)
            rows.append(
                ",".join(
                    csv_cell(value)
                    for value in (section.name, page.title, page.owner, section.page_type, attachments)
                )
            )
    (OUTPUT_DIR / "manifest.csv").write_text("\n".join(rows) + "\n", encoding="utf-8")


def write_readme() -> None:
    total_pages = sum(len(section.pages) for section in SECTIONS)
    total_attachments = sum(len(page.attachments) for section in SECTIONS for page in section.pages)
    lines = [
        f"# OneNote Content Pack - {COMPANY}",
        "",
        f"Generated pages: {total_pages}",
        f"Generated attachments: {total_attachments}",
        "",
        "A generic software-company knowledge base for testing the RAG assistant. Each page",
        "carries concrete, answerable facts (specific numbers, owners, ports, policies) so the",
        "chat returns real answers, not generic prose.",
        "",
        "## Sections",
        "",
    ]
    for index, section in enumerate(SECTIONS, start=1):
        lines.append(f"{index}. **{section.name}** - {SECTION_GUIDANCE[section.name]}")
    lines += [
        "",
        "## Import into OneNote",
        "",
        "1. Create one OneNote section per folder (drop the numeric prefix).",
        "2. For each `.html` file, open it in a browser, copy the rendered page, and paste it into a new OneNote page.",
        "3. For files named `<page>__<attachment>.<ext>`, attach them to that same OneNote page (Insert > File Attachment).",
        "4. Reindex so the app sees the new content:",
        "",
        "```powershell",
        "docker compose run --rm --build sync-worker onenote_bootstrap",
        "```",
        "",
        "## Attachments",
        "",
        "Attachments are real `.md`, `.txt`, `.docx`, `.pptx`, and `.pdf` files. Each contains at",
        "least one fact not repeated verbatim in the page body, so you can verify that attachment",
        "content is retrieved (e.g. ask about a value that only appears in the attached file).",
    ]
    (OUTPUT_DIR / "README.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


# --------------------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------------------

def render_list(items) -> str:
    return "\n      ".join(f"<li>{escape(item)}</li>" for item in items)


def comma_list(values) -> str:
    return ", ".join(values)


def safe_name(value: str) -> str:
    normalized = re.sub(r"[^A-Za-z0-9]+", "_", value).strip("_")
    return normalized or "page"


def csv_cell(value: str) -> str:
    return '"' + value.replace('"', '""') + '"'


if __name__ == "__main__":
    main()
